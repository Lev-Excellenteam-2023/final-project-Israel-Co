from pptx import Presentation
import re
from typing import List, Tuple


class PresentationParser:
    """
    Parsed presentation. Extracting each slide text and gathering them by index.
    """

    def __init__(self, path: str):
        """
        constructor
        :param path: Path of presentation file
        """
        self._slide_counter = 0
        self._slides = {}
        self._parse_presentation(path)

    def _parse_presentation(self, path: str):
        """
        Parse the presentation. Extract the text from each slide and save it by subject.
        :param path: Path of presentation file
        :return: Parsed presentation divided into a dictionary, the Key is slide number and the value is the text of the slide.
        """
        presentation = Presentation(path)
        for slide in presentation.slides:
            slide_text = self._extract_text_from_slide(slide)
            self._slides[self._slide_counter] = slide_text
            self._slide_counter += 1

    def _extract_text_from_slide(self, slide) -> str:
        """
        Extract text from the slide
        :param slide: Slide to extract text from
        :return: Text of the slide
        """
        slide_text_list = [shape.text + '\n' for shape in slide.shapes if hasattr(shape, 'text')]
        slide_text = ''.join(slide_text_list)

        return clean_weird_whitespaces(slide_text)

    def get_slide(self) -> Tuple[int, str]:
        """
        Generator which returns list of tuples. Each tuple contains slide number and the text slide.
        :return: Generator that gives tuple of slide number and the text slide.
        """
        for current_slide in self._slides:
            yield current_slide, self._slides[current_slide]

    def get_num_of_slides(self) -> int:
        return self._slide_counter


def clean_weird_whitespaces(text: str) -> str:
    """
    Clean weird whitespaces likewise two spaces, tabs est.
    :param text: Text to clean
    :return: The cleaned text
    """
    text = re.sub('^\n+', '', text)
    text = re.sub('\n\n+', '\n', text)
    text = re.sub('  +', ' ', text)
    text = re.sub('\t\t+', '\t', text)
    return text
