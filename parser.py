from pptx import Presentation
import re


class PresentationParser():
    def __init__(self):
        self._slide_counter = 1
        self._presentation_subjects = {}

    def parse_presentation(self, path):
        presentation = Presentation(path)
        for slide in presentation.slides:
            slide_title = slide.shapes.title.text
            slide_text = "".join(shape.text for shape in slide.shapes if hasattr(shape, 'text'))
            slide_text = clean_weird_whitespaces(slide_text)

            if slide_title != slide_text:
                if slide_title in self._presentation_subjects:
                    self._presentation_subjects[slide_title] += [(self._slide_counter, slide_text)]
                else:
                    self._presentation_subjects[slide_title] = [(self._slide_counter, slide_text)]

            self._slide_counter += 1


def clean_weird_whitespaces(text: str):
    text = re.sub('\n\n+', '\n', text)
    text = re.sub('  +', ' ', text)
    text = re.sub('\t\t+', '\t', text)
    return text


p = PresentationParser()
p.parse_presentation(r'C:\Users\josh5\Downloads\git.pptx')
print(p._presentation_subjects)
